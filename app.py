"""
StudySnap AI - Flask Application
Production-ready deployment for Render
"""

# CRITICAL: Load environment variables FIRST before any other imports
from dotenv import load_dotenv
load_dotenv()

import os
import re
import uuid
import threading
import time
import tempfile
from functools import wraps

from flask import Flask, render_template, request, session, jsonify, redirect, url_for
from youtube_transcript_api import YouTubeTranscriptApi

# Validate required environment variables before importing services
required_env_vars = ['GEMINI_API_KEY', 'MONGO_URI']
missing_vars = [var for var in required_env_vars if not os.getenv(var)]
if missing_vars:
    raise RuntimeError(f"Missing required environment variables: {', '.join(missing_vars)}")

# Import services (after env vars are validated)
from services.summary_gen import generate_study_material, _has_devanagari, _translate_to_english
from services.ai_chat import chat_with_ai
from services.mock_ai import generate_mock_test
from services.sample_paper_ai import generate_sample_paper
from services.study_pack_ai import generate_study_pack
from services.testdb import (
    create_user, verify_user, save_history, get_history,
    save_study_data, load_study_data, find_study_data_by_url, study_data_collection
)

# ── Flask App Configuration ──────────────────────────────────────────────────
app = Flask(
    __name__,
    template_folder="templates",
    static_folder="static"
)
app.secret_key = os.getenv('SECRET_KEY', 'studysnap-secret-key-change-in-production')

# ── Background Job Store ──────────────────────────────────────────────────────
_jobs = {}
_jobs_lock = threading.Lock()

# ── Study Data In-Memory Cache ────────────────────────────────────────────────
_study_data_cache = {}          # data_key -> {'data': ..., 'video_url': ..., 'ts': ...}
_study_data_cache_lock = threading.Lock()
_CACHE_TTL = 1800               # 30 minutes

# Tracks keys whose Hindi has already been fixed (avoid re-translating)
_translated_keys = set()
_translated_keys_lock = threading.Lock()


def _cache_store(key, data, video_url):
    """Write a study data entry into the in-memory cache."""
    with _study_data_cache_lock:
        _study_data_cache[key] = {'data': data, 'video_url': video_url, 'ts': time.time()}


def _cached_load_study_data(key):
    """Return study data from cache; fall back to MongoDB only on cache miss."""
    if not key:
        return None, ''
    now = time.time()
    with _study_data_cache_lock:
        entry = _study_data_cache.get(key)
        if entry and now - entry['ts'] < _CACHE_TTL:
            return entry['data'], entry['video_url']
    # Cache miss — fetch from DB and store
    try:
        data, video_url = load_study_data(key)
        if data is not None:
            _cache_store(key, data, video_url)
        return data, video_url
    except Exception as e:
        print(f'[cache] Error loading study data: {e}')
        return None, ''


def _create_job(job_id):
    with _jobs_lock:
        _jobs[job_id] = {'status': 'pending', 'result': None, 'error': None, 'ts': time.time()}


def _set_job_done(job_id, result):
    with _jobs_lock:
        if job_id in _jobs:
            _jobs[job_id].update({'status': 'done', 'result': result, 'error': None})


def _set_job_error(job_id, error):
    with _jobs_lock:
        if job_id in _jobs:
            _jobs[job_id].update({'status': 'error', 'result': None, 'error': error})


def _get_job(job_id):
    with _jobs_lock:
        # Purge stale jobs older than 45 minutes to free memory
        cutoff = time.time() - 2700
        stale = [k for k, v in _jobs.items() if v['ts'] < cutoff]
        for k in stale:
            del _jobs[k]
        return dict(_jobs.get(job_id, {}))


# ── Helper Functions ──────────────────────────────────────────────────────────

def extract_video_id(url):
    """Extract YouTube video ID from various URL formats."""
    patterns = [
        r'(?:v=)([0-9A-Za-z_-]{11})',
        r'(?:youtu\.be\/)([0-9A-Za-z_-]{11})',
        r'(?:embed\/)([0-9A-Za-z_-]{11})',
        r'(?:shorts\/)([0-9A-Za-z_-]{11})',
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None


def _extract_text_from_transcript_entries(entries):
    """Normalize transcript entries (dict/object/string) into plain text."""
    if not entries:
        return ''
    parts = []
    for entry in entries:
        if isinstance(entry, dict):
            text = entry.get('text', '')
        elif hasattr(entry, 'text'):
            text = getattr(entry, 'text', '')
        else:
            text = str(entry)
        text = (text or '').strip()
        if text:
            parts.append(text)
    return ' '.join(parts).strip()


def _extract_text_from_transcript_object(transcript_obj):
    """Normalize transcript object/list into plain text."""
    if transcript_obj is None:
        return ''
    if isinstance(transcript_obj, str):
        return transcript_obj.strip()
    if isinstance(transcript_obj, list):
        return _extract_text_from_transcript_entries(transcript_obj)
    if hasattr(transcript_obj, 'fetch'):
        try:
            fetched = transcript_obj.fetch()
            return _extract_text_from_transcript_entries(fetched)
        except Exception:
            return ''
    return ''


def get_transcript(video_id):
    """Get transcript text from a YouTube video. Supports multiple API versions."""
    languages = ['en', 'hi', 'en-US', 'en-GB', 'en-IN']

    # Try the static method approach (version 0.6.x and 1.x)
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        text = _extract_text_from_transcript_object(transcript)
        if text:
            return text
    except Exception as e:
        print(f'[transcript] Primary method failed: {e}')

    # Try with language fallback
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=languages)
        text = _extract_text_from_transcript_object(transcript)
        if text:
            return text
    except Exception as e:
        print(f'[transcript] Language fallback failed: {e}')

    # Try transcript listing APIs (covers more video/language cases)
    try:
        ytt = YouTubeTranscriptApi()
        list_fn = None
        if hasattr(ytt, 'list'):
            list_fn = ytt.list
        elif hasattr(YouTubeTranscriptApi, 'list_transcripts'):
            list_fn = YouTubeTranscriptApi.list_transcripts
        if list_fn:
            transcript_list = list_fn(video_id)
            for lang in languages:
                try:
                    if hasattr(transcript_list, 'find_transcript'):
                        chosen = transcript_list.find_transcript([lang])
                        text = _extract_text_from_transcript_object(chosen)
                        if text:
                            return text
                except Exception:
                    continue
            for finder in ('find_manually_created_transcript', 'find_generated_transcript'):
                fn = getattr(transcript_list, finder, None)
                if not fn:
                    continue
                try:
                    chosen = fn(languages)
                    text = _extract_text_from_transcript_object(chosen)
                    if text:
                        return text
                except Exception:
                    continue
    except Exception as e:
        print(f'[transcript] Transcript-list fallback failed: {e}')

    # Try instance-based approach (newer versions)
    try:
        ytt = YouTubeTranscriptApi()
        if hasattr(ytt, 'fetch'):
            transcript = ytt.fetch(video_id)
            text = _extract_text_from_transcript_object(transcript)
            if text:
                return text
        elif hasattr(ytt, 'get_transcript'):
            transcript = ytt.get_transcript(video_id)
            text = _extract_text_from_transcript_object(transcript)
            if text:
                return text
    except Exception as e:
        print(f'[transcript] Instance method failed: {e}')

    return None


# ── Auth helper ──────────────────────────────────────────────────────────────

def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get('logged_in'):
            return redirect(url_for('signin'))
        return f(*args, **kwargs)
    return decorated


# ── Routes ────────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    if session.get('logged_in'):
        return redirect(url_for('home'))
    return redirect(url_for('signin'))


@app.route("/home")
@login_required
def home():
    return render_template("home.html")


@app.route("/signin", methods=["GET", "POST"])
def signin():
    if request.method == "POST":
        email = request.form.get("email", "").strip()
        password = request.form.get("password", "")
        if not email or not password:
            return render_template("signin.html", error="Please fill in all fields.")
        try:
            ok, result = verify_user(email, password)
            if ok:
                session['logged_in'] = True
                session['user_email'] = result['email']
                session['user_name'] = result['name']
                return redirect(url_for('home'))
            return render_template("signin.html", error=result)
        except Exception as e:
            print(f'[signin] Database error: {e}')
            return render_template("signin.html", error="Database connection failed. Please try again.")
    return render_template("signin.html")


@app.route("/signup", methods=["GET", "POST"])
def signup():
    if request.method == "POST":
        name = request.form.get("name", "").strip()
        email = request.form.get("email", "").strip()
        password = request.form.get("password", "")
        confirm = request.form.get("confirm_password", "")
        if not name or not email or not password:
            return render_template("signup.html", error="Please fill in all fields.")
        if password != confirm:
            return render_template("signup.html", error="Passwords do not match.")
        if len(password) < 6:
            return render_template("signup.html", error="Password must be at least 6 characters.")
        try:
            ok, msg = create_user(name, email, password)
            if ok:
                session['logged_in'] = True
                session['user_email'] = email
                session['user_name'] = name
                return redirect(url_for('home'))
            return render_template("signup.html", error=msg)
        except Exception as e:
            print(f'[signup] Database error: {e}')
            return render_template("signup.html", error="Database connection failed. Please try again.")
    return render_template("signup.html")


@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for('signin'))


@app.route("/analize")
@login_required
def analize():
    return render_template("analize.html")


@app.route("/summary")
@login_required
def summary():
    key = session.get('data_key')
    data, video_url = _cached_load_study_data(key)
    if data and isinstance(data, dict):
        data = _normalize_study_data(data, key, video_url)
    return render_template("summary.html", data=data, video_url=video_url)


def _normalize_study_data(data, data_key=None, video_url=''):
    """Ensure data dict has all keys the summary template expects.
    Translates Hindi in English fields once and persists the fix."""
    data.setdefault('video_title', 'Untitled')
    data.setdefault('chapter_title', '')
    data.setdefault('main_summary', {})
    data.setdefault('conclusion', '')
    if not isinstance(data.get('key_terms'), list):
        data['key_terms'] = []
    if not isinstance(data.get('core_observations'), list):
        data['core_observations'] = []
    if not isinstance(data.get('highlights'), list):
        data['highlights'] = []
    for key in ('important_points', 'priority_topics', 'hindi_notes', 'main_summary'):
        val = data.get(key)
        if not isinstance(val, dict):
            data[key] = {'short': str(val) if val else '', 'medium': str(val) if val else '', 'detailed': str(val) if val else ''}
        else:
            for level in ('short', 'medium', 'detailed'):
                if not isinstance(val.get(level), str):
                    val[level] = str(val.get(level, ''))
    # priority_topics_all: always the full comprehensive content, not filtered by detail level
    pt = data.get('priority_topics', {})
    if isinstance(pt, dict):
        data['priority_topics_all'] = pt.get('detailed') or pt.get('medium') or pt.get('short') or ''
    elif isinstance(pt, str):
        data['priority_topics_all'] = pt
    else:
        data['priority_topics_all'] = ''

    # ── Translate Hindi leaked into English fields (one-time, then persist) ──
    if data_key:
        with _translated_keys_lock:
            already_done = data_key in _translated_keys
        if not already_done:
            changed = False
            for field in ('important_points', 'main_summary'):
                val = data.get(field, {})
                if isinstance(val, dict):
                    for level in ('short', 'medium', 'detailed'):
                        text = val.get(level, '')
                        if text and _has_devanagari(text):
                            try:
                                val[level] = _translate_to_english(text)
                                changed = True
                            except Exception as e:
                                print(f'[normalize] Translation failed: {e}')
            # priority_topics dict levels
            ptd = data.get('priority_topics', {})
            if isinstance(ptd, dict):
                for level in ('short', 'medium', 'detailed'):
                    text = ptd.get(level, '')
                    if text and _has_devanagari(text):
                        try:
                            ptd[level] = _translate_to_english(text)
                            changed = True
                        except Exception as e:
                            print(f'[normalize] Translation failed: {e}')
            # Refresh priority_topics_all after possible translation
            if changed:
                pt2 = data.get('priority_topics', {})
                if isinstance(pt2, dict):
                    data['priority_topics_all'] = pt2.get('detailed') or pt2.get('medium') or pt2.get('short') or ''
                # Persist: update cache + DB so this never runs again
                _cache_store(data_key, data, video_url)
                try:
                    study_data_collection.update_one({'key': data_key}, {'$set': {'data': data}})
                except Exception as e:
                    print(f'[normalize] DB persist failed: {e}')
            with _translated_keys_lock:
                _translated_keys.add(data_key)

    return data


@app.route("/mock")
@login_required
def mock():
    key = session.get('data_key')
    data, _ = _cached_load_study_data(key)
    has_summary = data is not None
    return render_template("mock.html", has_summary=has_summary, data_key=key or '')


@app.route("/pdf")
@login_required
def pdf():
    key = session.get('data_key')
    data, _ = _cached_load_study_data(key)
    has_summary = data is not None
    return render_template("pdf.html", has_summary=has_summary)


@app.route("/ai")
@login_required
def ai():
    return render_template("ai.html")


@app.route("/settings")
@login_required
def settings():
    return render_template("settings.html")


@app.route("/history")
@login_required
def history():
    email = session.get('user_email')
    try:
        entries = get_history(email) if email else []
    except Exception as e:
        print(f'[history] Database error: {e}')
        entries = []
    return render_template("history.html", entries=entries)


@app.route("/history/view/<data_key>")
@login_required
def history_view(data_key):
    data, video_url = _cached_load_study_data(data_key)
    if data is None:
        return redirect(url_for('history'))
    session['data_key'] = data_key
    return redirect(url_for('summary'))


@app.route("/sample-paper")
@login_required
def sample_paper():
    key = session.get('data_key')
    data, _ = _cached_load_study_data(key)
    has_summary = data is not None
    return render_template("samplePaper.html", has_summary=has_summary, data_key=key or '')


# ── File text extraction helpers ──────────────────────────────────────────────

_uploaded_texts = {}  # file_key -> extracted text string
_uploaded_texts_lock = threading.Lock()


def _extract_pdf_text(filepath):
    """Extract text from a PDF file."""
    from pypdf import PdfReader
    reader = PdfReader(filepath)
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return '\n'.join(parts)


def _extract_pptx_text(filepath):
    """Extract text from a PPTX file."""
    from pptx import Presentation
    prs = Presentation(filepath)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
    return '\n'.join(parts)


# ── API Endpoints ─────────────────────────────────────────────────────────────

@app.route("/api/upload-file", methods=["POST"])
def api_upload_file():
    """Upload a PDF/PPTX file, extract its text, and store it for generation."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded.'}), 400
    file = request.files['file']
    if not file.filename:
        return jsonify({'error': 'No file selected.'}), 400

    filename = file.filename
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    if ext not in ('pdf', 'pptx', 'ppt'):
        return jsonify({'error': 'Unsupported file type. Please upload a PDF or PPTX file.'}), 400

    # Save to a temp file for processing
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.' + ext) as tmp:
            file.save(tmp.name)
            tmp_path = tmp.name

        if ext == 'pdf':
            text = _extract_pdf_text(tmp_path)
        else:
            text = _extract_pptx_text(tmp_path)

        os.unlink(tmp_path)
    except Exception as e:
        return jsonify({'error': f'Failed to read file: {str(e)}'}), 400

    if not text or len(text.strip()) < 50:
        return jsonify({'error': 'Could not extract enough text from the file. The file may be image-based or empty.'}), 400

    file_key = str(uuid.uuid4())
    with _uploaded_texts_lock:
        # Purge old entries (keep max 100)
        if len(_uploaded_texts) > 100:
            oldest = list(_uploaded_texts.keys())[:50]
            for k in oldest:
                del _uploaded_texts[k]
        _uploaded_texts[file_key] = text

    return jsonify({'success': True, 'file_key': file_key})


@app.route("/api/generate-from-file", methods=["POST"])
def api_generate_from_file():
    """Generate study material from an uploaded PDF/PPTX file."""
    data = request.get_json()
    file_key = data.get('file_key', '')
    exam_mode = data.get('exam_mode', 'College/University')
    language = data.get('language', 'English (Pure)')
    difficulty = data.get('difficulty', 'Medium (Standard)')
    study_depth = data.get('study_depth', '40')

    with _uploaded_texts_lock:
        transcript = _uploaded_texts.pop(file_key, None)

    if not transcript:
        return jsonify({'error': 'File not found or expired. Please upload again.'}), 400

    try:
        summary_data = generate_study_material(transcript, exam_mode, language, difficulty, study_depth)
        source_label = 'Uploaded Document'
        key = save_study_data(summary_data, source_label)
        _cache_store(key, summary_data, source_label)
        session['data_key'] = key
        email = session.get('user_email')
        if email:
            title = summary_data.get('video_title', 'Untitled') if isinstance(summary_data, dict) else 'Untitled'
            save_history(email, title, source_label, key)
        return jsonify({'success': True})
    except Exception as e:
        print(f'[summary] Exception ({type(e).__name__}): {e}')
        return jsonify({'error': f'AI generation failed: {str(e)}. Please try again.'}), 500


@app.route("/api/generate", methods=["POST"])
def api_generate():
    """Process YouTube video and generate AI study material."""
    data = request.get_json()
    url = data.get('url', '')
    exam_mode = data.get('exam_mode', 'College/University')
    language = data.get('language', 'English (Pure)')
    difficulty = data.get('difficulty', 'Medium (Standard)')
    study_depth = data.get('study_depth', '40')

    # Extract video ID
    video_id = extract_video_id(url)
    if not video_id:
        return jsonify({'error': 'Invalid YouTube URL. Please paste a valid YouTube video link.'}), 400

    # Get transcript
    transcript = get_transcript(video_id)
    if not transcript:
        return jsonify({'error': 'Could not extract transcript. The video may not have captions/subtitles enabled.'}), 400

    # Check cache — if we already generated notes for this URL, reuse them
    try:
        cached_data, cached_key = find_study_data_by_url(url)
        if cached_data and cached_key:
            session['data_key'] = cached_key
            _cache_store(cached_key, cached_data, url)
            email = session.get('user_email')
            if email:
                title = cached_data.get('video_title', 'Untitled') if isinstance(cached_data, dict) else 'Untitled'
                save_history(email, title, url, cached_key)
            return jsonify({'success': True, 'cached': True})
    except Exception as e:
        print(f'[generate] Cache lookup failed: {e}')

    # Generate AI content
    try:
        summary_data = generate_study_material(transcript, exam_mode, language, difficulty, study_depth)
        key = save_study_data(summary_data, url)
        _cache_store(key, summary_data, url)
        session['data_key'] = key
        # Save to user history
        email = session.get('user_email')
        if email:
            title = summary_data.get('video_title', 'Untitled') if isinstance(summary_data, dict) else 'Untitled'
            save_history(email, title, url, key)
        return jsonify({'success': True})
    except Exception as e:
        print(f'[summary] Exception ({type(e).__name__}): {e}')
        return jsonify({'error': f'AI generation failed: {str(e)}. Please try again.'}), 500


@app.route("/api/chat", methods=["POST"])
def api_chat():
    """Ask AI chat endpoint."""
    data = request.get_json()
    user_message = data.get('message', '').strip()
    history = data.get('history', [])

    if not user_message:
        return jsonify({'error': 'Empty message'}), 400

    # Load summary context if available
    summary_context = None
    key = session.get('data_key')
    if key:
        summary_context, _ = _cached_load_study_data(key)

    try:
        reply = chat_with_ai(user_message, history, summary_context)
        return jsonify({'reply': reply})
    except Exception as e:
        print(f'[chat] AI error: {e}')
        return jsonify({'error': f'AI chat failed: {str(e)}'}), 500


@app.route("/api/mock-generate", methods=["POST"])
def api_mock_generate():
    """Generate mock test from summary data (background job)."""
    key = session.get('data_key')
    data, _ = _cached_load_study_data(key)
    if not data:
        return jsonify({'error': 'No summary available. Please generate a summary first from the home page.'}), 400
    job_id = str(uuid.uuid4())
    _create_job(job_id)
    def _run(d):
        try:
            _set_job_done(job_id, generate_mock_test(d))
        except Exception as e:
            _set_job_error(job_id, f'Failed to generate mock test: {str(e)}')
    threading.Thread(target=_run, args=(data,), daemon=True).start()
    return jsonify({'success': True, 'job_id': job_id})


@app.route("/api/sample-paper-generate", methods=["POST"])
def api_sample_paper_generate():
    """Generate sample paper from summary data (background job)."""
    key = session.get('data_key')
    data, _ = _cached_load_study_data(key)
    if not data:
        return jsonify({'error': 'No summary available. Please generate a summary first from the home page.'}), 400
    body = request.get_json(silent=True) or {}
    counts = {
        'mcq':     max(1, min(20, int(body.get('mcq',     5)))),
        'fillups': max(1, min(15, int(body.get('fillups', 5)))),
        'short':   max(1, min(15, int(body.get('short',   5)))),
        'long':    max(1, min(10, int(body.get('long',    5)))),
    }
    job_id = str(uuid.uuid4())
    _create_job(job_id)
    def _run(d, c):
        try:
            _set_job_done(job_id, generate_sample_paper(d, c))
        except Exception as e:
            _set_job_error(job_id, f'Failed to generate sample paper: {str(e)}')
    threading.Thread(target=_run, args=(data, counts), daemon=True).start()
    return jsonify({'success': True, 'job_id': job_id})


@app.route("/api/study-pack-generate", methods=["POST"])
def api_study_pack_generate():
    """Generate comprehensive study pack for PDF download (background job)."""
    key = session.get('data_key')
    data, _ = _cached_load_study_data(key)
    if not data:
        return jsonify({'error': 'No summary available. Please generate a summary first from the home page.'}), 400
    job_id = str(uuid.uuid4())
    _create_job(job_id)
    def _run(d):
        try:
            _set_job_done(job_id, generate_study_pack(d))
        except Exception as e:
            print(f'[study-pack] Exception: {e}')
            _set_job_error(job_id, f'Failed to generate study pack: {str(e)}')
    threading.Thread(target=_run, args=(data,), daemon=True).start()
    return jsonify({'success': True, 'job_id': job_id})


@app.route("/api/job/<job_id>", methods=["GET"])
def api_job_status(job_id):
    """Poll a background job's status and result."""
    job = _get_job(job_id)
    if not job:
        return jsonify({'status': 'not_found'}), 404
    # Don't send full result in every poll – only when done
    response = {'status': job['status'], 'error': job.get('error')}
    if job['status'] == 'done':
        response['data'] = job['result']
    return jsonify(response)


# ── Health Check Endpoint ─────────────────────────────────────────────────────

@app.route("/health")
def health_check():
    """Health check endpoint for Render."""
    return jsonify({'status': 'healthy'}), 200


# ── Application Entry Point ───────────────────────────────────────────────────

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
